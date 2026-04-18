"""
deck_generator.py — Bespoke 6-Slide PPTX Pitch Deck Generator
==============================================================
Generates company-specific pitch decks using python-pptx.

Design: dark 0A0F1C background, card-based layout, rotating accents.
Copy: data-driven from enrichment; AI-enhanced when ANTHROPIC_API_KEY is set.

Slide structure:
  1  Cover              — headline names the specific problem
  2  The Problem        — 2x2 card grid + tension line
  3  Powered by Claude  — numbered steps + fake terminal
  4  Market Reality     — named competitor cards + insight box
  5  The System         — 3 outcome sentences only
  6  Pilot + CTA        — metrics + booking panel

Usage:
    from deck_generator import generate_deck
    path = generate_deck(prospect_dict)
"""

import os
import re
import json
import shutil
import subprocess
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from database import DB_PATH, get_all_prospects
from settings import get_calendar_link

OUTPUT_DIR = "decks"

# Slide dimensions — 16:9 widescreen
SW = Inches(13.33)
SH = Inches(7.5)

# ---------------------------------------------------------------------------
# Palette (no # prefix — RGBColor needs raw hex)
# ---------------------------------------------------------------------------
BG       = "0A0F1C"
CARD_BG  = "131929"
CARD_BD  = "3D4A5C"
TEXT_PRI = "E2E8F0"
TEXT_SEC = "94A3B8"
TEXT_MUT = "64748B"
TERM_BG  = "0D1117"
TERM_GRN = "4ADE80"
TERM_CMD = "7DD3FC"
ACCENTS  = ["7C3AED", "F97316", "06B6D4", "10B981"]
DANGER   = "EF4444"
INDIGO   = "6366F1"

FOOTER_TXT = "Antigravity  \u00b7  Confidential"
CTA_LINK   = get_calendar_link()

# ---------------------------------------------------------------------------
# Visual themes + layout styles
# ---------------------------------------------------------------------------
# Each theme entry: palette keys (swap color globals) + "STYLE" sub-dict.
# STYLE controls layout/sizing decisions read by slide functions.
# ---------------------------------------------------------------------------

# Module-level STYLE dict — overwritten by _apply_theme context manager
STYLE: dict = {}   # populated below after THEMES definition

THEMES = {
    # ── 1  Dark Indigo ──────────────────────────────────────────────────────
    # Deep navy · purple/orange/cyan · left-hero cover · stacked stat cards
    "dark_indigo": {
        "BG":       "0A0F1C",
        "CARD_BG":  "131929",
        "CARD_BD":  "3D4A5C",
        "TEXT_PRI": "E2E8F0",
        "TEXT_SEC": "94A3B8",
        "TEXT_MUT": "64748B",
        "TERM_BG":  "0D1117",
        "TERM_GRN": "4ADE80",
        "TERM_CMD": "7DD3FC",
        "ACCENTS":  ["7C3AED", "F97316", "06B6D4", "10B981"],
        "STYLE": {
            "accent_bar":       "top",
            "accent_bar_h":     0.055,
            "h1": 44, "h2": 34, "h3": 16, "body": 11,
            "card_border_pt":   0.75,
            "card_stripe":      "dot",       # tiny accent dot on problem cards
            "section_prefix":   "",
            "section_suffix":   "",
            "section_mono":     False,
            "cover_layout":     "hero_left",
            "slide5_layout":    "stripe_cards",
        },
    },
    # ── 2  Charcoal Gold ────────────────────────────────────────────────────
    # Pure black · gold accents · centered cover · borderless · COMPACT 16:9
    "charcoal_gold": {
        "BG":       "0C0C0C",
        "CARD_BG":  "181818",
        "CARD_BD":  "3A3020",
        "TEXT_PRI": "F5F0E8",
        "TEXT_SEC": "B8A98A",
        "TEXT_MUT": "6B5F4A",
        "TERM_BG":  "111008",
        "TERM_GRN": "D4A853",
        "TERM_CMD": "E8C664",
        "ACCENTS":  ["D4A853", "C4973A", "E8C664", "F0B429"],
        "STYLE": {
            "accent_bar":       "left",      # vertical left-side stripe
            "accent_bar_h":     0.10,        # width of left stripe
            "h1": 50, "h2": 38, "h3": 17, "body": 11,
            "card_border_pt":   0,           # borderless cards
            "card_stripe":      "top",       # top stripe on each card
            "section_prefix":   "— ",
            "section_suffix":   "",
            "section_mono":     False,
            "cover_layout":     "centered_bottom",
            "slide5_layout":    "stripe_cards",
            "slide_scale":      0.75,   # compact 16:9  (10 × 5.625 in)
        },
    },
    # ── 3  Midnight Teal ────────────────────────────────────────────────────
    # Dark teal · cyan/emerald · reversed split cover · heavy borders
    "midnight_teal": {
        "BG":       "061616",
        "CARD_BG":  "0C2222",
        "CARD_BD":  "1A5050",
        "TEXT_PRI": "E0F7F7",
        "TEXT_SEC": "7ABFBF",
        "TEXT_MUT": "3D7070",
        "TERM_BG":  "040F0F",
        "TERM_GRN": "2DD4BF",
        "TERM_CMD": "67E8F9",
        "ACCENTS":  ["14B8A6", "06B6D4", "34D399", "22D3EE"],
        "STYLE": {
            "accent_bar":       "bottom",    # accent line at foot of slide
            "accent_bar_h":     0.055,
            "h1": 40, "h2": 30, "h3": 15, "body": 10.5,
            "card_border_pt":   1.5,         # thick visible border
            "card_stripe":      "none",
            "section_prefix":   "[ ",
            "section_suffix":   " ]",
            "section_mono":     True,
            "cover_layout":     "reversed_split",
            "slide5_layout":    "pipe_lines",  # no cards — just large pipe-separated lines
        },
    },
    # ── 4  Dark Slate Red ───────────────────────────────────────────────────
    # Near-black · bold red/amber · thick bar · oversized type · fullbleed cover
    "dark_slate_red": {
        "BG":       "0F0F0F",
        "CARD_BG":  "1A1414",
        "CARD_BD":  "3D2020",
        "TEXT_PRI": "F5F0F0",
        "TEXT_SEC": "B09090",
        "TEXT_MUT": "6B4848",
        "TERM_BG":  "0A0808",
        "TERM_GRN": "FB923C",
        "TERM_CMD": "FCD34D",
        "ACCENTS":  ["EF4444", "F97316", "FBBF24", "F43F5E"],
        "STYLE": {
            "accent_bar":       "top_thick",  # very thick top bar (0.18 in)
            "accent_bar_h":     0.18,
            "h1": 54, "h2": 40, "h3": 18, "body": 11,
            "card_border_pt":   0,
            "card_stripe":      "left",       # left accent stripe replaces border
            "section_prefix":   "",
            "section_suffix":   "",
            "section_mono":     False,
            "cover_layout":     "fullbleed_bold",
            "slide5_layout":    "stripe_cards",
            "slide_scale":      0.75,   # compact 16:9  (10 × 5.625 in)
        },
    },
}

# Seed the default STYLE so the module works without a theme context
STYLE = dict(THEMES["dark_indigo"]["STYLE"])

import contextlib

@contextlib.contextmanager
def _apply_theme(theme_name: str):
    """Temporarily swap module-level palette + style globals."""
    import deck_generator as _self
    theme = THEMES.get(theme_name, THEMES["dark_indigo"])
    style = theme.get("STYLE", THEMES["dark_indigo"]["STYLE"])
    palette_keys = [k for k in theme if k != "STYLE"]
    saved_palette   = {k: getattr(_self, k) for k in palette_keys}
    saved_style     = dict(_self.STYLE)
    saved_scale     = _self.SLIDE_SCALE
    try:
        for k, v in theme.items():
            if k != "STYLE":
                setattr(_self, k, v)
        _self.STYLE       = dict(style)
        _self.SLIDE_SCALE = float(style.get("slide_scale", 1.0))
        yield
    finally:
        for k, v in saved_palette.items():
            setattr(_self, k, v)
        _self.STYLE       = saved_style
        _self.SLIDE_SCALE = saved_scale

# ---------------------------------------------------------------------------
# Copy quality gates
# ---------------------------------------------------------------------------
BANNED_PHRASES = [
    "it appears", "based on available data", "may indicate", "could suggest",
    "this indicates", "this suggests", "likely dependent on", "appears to rely on",
    "likely could", "seemingly", "seems like", "might be", "could be",
    "it is worth noting", "we noticed", "it looks like",
    "in today's competitive landscape", "current acquisition model",
    "worth exploring", "leverage", "synergies", "unlock growth opportunities",
    "holistic", "robust framework", "end-to-end", "this approach has a ceiling",
]

TENSION_MARKERS = [
    # constraint / ceiling language
    "caps your", "caps growth", "this caps", "hard.", "ceiling",
    "stops when", "pipeline stops", "growth stops", "resets to zero",
    # absence / invisibility
    "you do not", "you don't", "they do not", "they don't",
    "never see", "invisible to", "not in it", "just not in",
    "does not exist", "they do not exist", "nothing is missed",
    # dependency / fragility
    "you are renting", "rented growth", "by default", "no backup",
    "depends on being found", "depends on", "no other way in",
    "not a system", "hiring does not fix",
    # timing / first-mover
    "already in your", "first contact", "immediately.",
    "before they go looking", "before they start", "before inbound",
    "before they search", "before they find",
    # risk / cost
    "if no results", "if the pipeline", "zero risk", "no charge",
    "cost if no", "no long-term", "no contract",
    # direct tension phrases
    "that matters more", "more than having", "sets the frame",
    "already happened", "without you", "without them",
]

# ---------------------------------------------------------------------------
# AI system prompt template
# ---------------------------------------------------------------------------
DECK_SYSTEM_PROMPT_TEMPLATE = """
ROLE

You are a senior sales deck designer and copywriter. You generate a bespoke 6-slide .pptx pitch deck for a specific prospect company using PptxGenJS. Every deck must feel hand-crafted for that company, not templated.

INPUTS YOU WILL RECEIVE

{{company_name}}
{{contact_name}}
{{industry}}
{{target_buyer}}
{{core_offer}}
{{current_acquisition}}
{{competitor_1}}
{{competitor_2}}
{{pain_point}}
{{date}}

DESIGN RULES - NEVER BREAK THESE

- Background 0A0F1C on every slide. No white. No light grey.
- Cards: 131929 fill, 3D4A5C border, drop shadow via makeShadow() factory.
- Accent colours rotate: 7C3AED, F97316, 06B6D4, 10B981
- Full-width accent bar 0.055 inches tall at top of every slide
- Section label: spaced caps, 9pt, accent colour, top left
- Footer every slide: "LeadGen AI · Confidential" left, page number right, 9pt grey
- No horizontal dividers. No bullet unicode characters. No # prefix on hex colours.
- margin: 0 on all text boxes. makeShadow must be a factory function called fresh each time.

SLIDE STRUCTURE - FILL WITH COMPANY DATA

Slide 1 - Cover
- Headline uses {{company_name}}'s specific problem
- "Built on Claude Code" badge top left
- 3 stat cards right side: qualified calls/month, pilot length, risk guarantee
- Bottom: "Prepared for {{contact_name}} at {{company_name}} · {{date}}"

Slide 2 - The Problem
- Headline names the specific constraint facing {{company_name}} based on {{current_acquisition}}
- 2x2 dark card grid referencing {{industry}} or {{target_buyer}}
- Tension line at the bottom

Slide 3 - Powered by Claude Code
- Left: 4-step numbered list written for {{industry}}
- Right: fake terminal card referencing {{company_name}}, {{target_buyer}}, and {{pain_point}}

Slide 4 - Market Reality
- {{competitor_1}} and {{competitor_2}} named explicitly in the two cards
- Body copy states specifically what they are doing in {{industry}}
- Insight box adapts "They get first contact. You do not." to the market context

Slide 5 - The System
- Three outcome lines only. No step lists.
- "We find the [{{target_buyer}}] already dealing with [{{pain_point}}]."
- "We reach them directly before they start searching for a solution."
- "We turn replies into booked conversations with your team."

Slide 6 - Pilot + CTA
- Headline: a short risk-reversal (e.g. "No results. No charge." or "Zero risk. Real pipeline.")
- Metric cards specific to {{industry}}
- CTA panel references {{company_name}} by name
- Pricing section: Pilot free, retainer from £1,500/month

COPY RULES - ENFORCED AT RENDER TIME

Every non-cover slide must end with or contain a short tension line — a single sentence that names a hard constraint, a missed opportunity, or an uncomfortable truth. Good tension lines are short and direct. Examples:
  "This caps your pipeline. Hard."
  "They get first contact. You do not."
  "First contact happens before inbound does."
  "Nothing is missed."
  "They do not exist to your best buyers yet."
Ban these phrases: end-to-end / leverage / holistic / robust / this suggests / this indicates / appears to rely on / current acquisition model / based on available data
The system slide must never be a process list. Always three outcome sentences.

QA STEPS - REQUIRED BEFORE DELIVERING

1. Convert to PDF with soffice
2. Rasterise with pdftoppm
3. Visually inspect for overlap, overflow, white backgrounds, missing company references, or generic copy
4. Run banned phrase check
5. Confirm every slide has at least one tension line
6. Fix issues and re-render before delivering

FINAL RULE

A founder at {{company_name}} should read this and think: "They actually understand our market."

Return ONLY valid JSON using this schema:
{
  "slide1": {
    "headline": "string",
    "subline": "string",
    "stats": [
      {"number": "4-8", "label": "qualified calls/month"},
      {"number": "30", "label": "day pilot"},
      {"number": "£0", "label": "if no results"}
    ],
    "prepared_for": "string",
    "tension": "string"
  },
  "slide2": {
    "headline": "string",
    "cards": [
      {"title": "string", "body": "string"},
      {"title": "string", "body": "string"},
      {"title": "string", "body": "string"},
      {"title": "string", "body": "string"}
    ],
    "tension": "string"
  },
  "slide3": {
    "headline": "string",
    "steps": ["string", "string", "string", "string"],
    "terminal_lines": ["string", "string", "string", "string", "string", "string"],
    "tension": "string"
  },
  "slide4": {
    "headline": "string",
    "comp1": {"name": "string", "line1": "string", "line2": "string"},
    "comp2": {"name": "string", "line1": "string", "line2": "string"},
    "insight": "string",
    "tension": "string"
  },
  "slide5": {
    "headline": "The System",
    "outcomes": ["string", "string", "string"],
    "tension": "string"
  },
  "slide6": {
    "headline": "No results. No charge.",
    "metrics": [
      {"number": "string", "label": "string"},
      {"number": "string", "label": "string"},
      {"number": "string", "label": "string"}
    ],
    "cta_body": "string",
    "pricing": "Pilot: Free  |  Ongoing: from £1,500/month  |  No long-term contract",
    "tension": "string"
  }
}
""".strip()


# ---------------------------------------------------------------------------
# Slide scale — overwritten by _apply_theme for compact (10×5.625) themes
# ---------------------------------------------------------------------------
SLIDE_SCALE: float = 1.0   # 1.0 = 13.33×7.5 in,  0.75 = 10×5.625 in (both 16:9)


def _I(n: float) -> int:
    """Scale-aware Inches() for slide drawing."""
    return Inches(n * SLIDE_SCALE)


def _P(n: float) -> int:
    """Scale-aware Pt() for font sizes and line widths."""
    return Pt(n * SLIDE_SCALE)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _fill_prompt_template(template: str, variables: dict) -> str:
    """Replace {{variables}} placeholders in the stored system prompt template."""
    filled = template
    for key, value in variables.items():
        filled = filled.replace(f"{{{{{key}}}}}", value or "")
    return filled


def build_deck_system_prompt(variables: dict) -> str:
    """Return the filled deck system prompt for a specific prospect."""
    return _fill_prompt_template(DECK_SYSTEM_PROMPT_TEMPLATE, variables)


def _blank_slide(prs: Presentation):
    """Add a blank-layout slide to the presentation."""
    blank = next(
        (l for l in prs.slide_layouts if l.name == 'Blank'),
        prs.slide_layouts[6],
    )
    return prs.slides.add_slide(blank)


def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(BG)


def _accent_bar(slide, accent: str):
    """Accent bar — position and size driven by STYLE."""
    pos  = STYLE.get("accent_bar", "top")
    h    = STYLE.get("accent_bar_h", 0.055)
    sw   = _I(13.33)   # scaled slide width
    sh   = _I(7.5)     # scaled slide height
    if pos in ("top", "top_thick"):
        bar = slide.shapes.add_shape(1, _I(0), _I(0), sw, _I(h))
    elif pos == "left":
        bar = slide.shapes.add_shape(1, _I(0), _I(0), _I(h), sh)
    elif pos == "bottom":
        bar = slide.shapes.add_shape(1, _I(0), sh - _I(h), sw, _I(h))
    else:
        bar = slide.shapes.add_shape(1, _I(0), _I(0), sw, _I(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()


def _footer(slide, page_num: int):
    """Footer: brand left, page number right, 7.5pt muted grey."""
    # Left
    tb = slide.shapes.add_textbox(_I(0.4), _I(7.22), _I(7), _I(0.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = FOOTER_TXT
    r.font.size = _P(7.5)
    r.font.color.rgb = _rgb(TEXT_MUT)

    # Right
    tb2 = slide.shapes.add_textbox(_I(11.5), _I(7.22), _I(1.43), _I(0.2))
    tf2 = tb2.text_frame
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_num)
    r2.font.size = _P(7.5)
    r2.font.color.rgb = _rgb(TEXT_MUT)


def _section_label(slide, text: str, accent: str, y: float = 0.13):
    """Spaced-caps section label — prefix/suffix from STYLE."""
    prefix = STYLE.get("section_prefix", "")
    suffix = STYLE.get("section_suffix", "")
    mono   = STYLE.get("section_mono", False)
    x_off  = 0.55 if STYLE.get("accent_bar") == "left" else 0.45
    tb = slide.shapes.add_textbox(_I(x_off), _I(y), _I(7), _I(0.22))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{prefix}{text.upper()}{suffix}"
    r.font.size = _P(8)
    r.font.bold = True
    r.font.color.rgb = _rgb(accent)
    if mono:
        r.font.name = "Courier New"
    else:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', '150')


def _shadow(shape):
    """
    Shadow factory — creates a fresh drop-shadow XML block every call.
    Never reuses the same object across shapes.
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = OxmlElement('a:effectLst')
    outerShdw = OxmlElement('a:outerShdw')
    outerShdw.set('blurRad', '50800')    # ~4pt
    outerShdw.set('dist',    '38100')    # ~3pt
    outerShdw.set('dir',     '2700000')  # 45° (down-right)
    outerShdw.set('rotWithShape', '0')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', '000000')
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '50000')
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def _card(slide, x: float, y: float, w: float, h: float, with_shadow: bool = True):
    """Dark card rectangle — border driven by STYLE."""
    shape = slide.shapes.add_shape(1, _I(x), _I(y), _I(w), _I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(CARD_BG)
    border_pt = STYLE.get("card_border_pt", 0.75)
    if border_pt > 0:
        shape.line.color.rgb = _rgb(CARD_BD)
        shape.line.width = _P(border_pt)
    else:
        shape.line.fill.background()
    if with_shadow:
        _shadow(shape)
    return shape


def _card_stripe(slide, x: float, y: float, w: float, h: float, accent: str):
    """Draw a thin accent stripe on a card (top or left) per STYLE."""
    mode = STYLE.get("card_stripe", "dot")
    if mode == "top":
        stripe = slide.shapes.add_shape(1, _I(x), _I(y), _I(w), _I(0.065))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _rgb(accent)
        stripe.line.fill.background()
    elif mode == "left":
        stripe = slide.shapes.add_shape(1, _I(x), _I(y), _I(0.07), _I(h))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _rgb(accent)
        stripe.line.fill.background()


def _tx(slide, x: float, y: float, w: float, h: float) -> object:
    """Add a transparent textbox. Returns text_frame with zero margins."""
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    return tf


def _run(tf, text: str, size: float, color: str,
         bold: bool = False, italic: bool = False,
         mono: bool = False, new_para: bool = True,
         align=PP_ALIGN.LEFT, space_before: float = 0) -> object:
    """Add a paragraph+run to a text_frame. First call uses existing para."""
    if new_para and len(tf.paragraphs) > 0 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = _P(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = _P(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = _rgb(color)
    if mono:
        r.font.name = 'Courier New'
    return p


# ---------------------------------------------------------------------------
# Copy quality validation
# ---------------------------------------------------------------------------

def _validate_copy(copy: dict) -> list:
    """Return list of validation errors across all slide copy."""
    errors = []
    full_text = json.dumps(copy).lower()

    banned = [p for p in BANNED_PHRASES if p in full_text]
    if banned:
        errors.append(f"Banned phrases found: {', '.join(banned)}")

    # Each non-cover slide must have at least one tension marker
    for key in ('slide2', 'slide3', 'slide4', 'slide5', 'slide6'):
        slide_text = json.dumps(copy.get(key, {})).lower()
        if not any(m in slide_text for m in TENSION_MARKERS):
            errors.append(f"{key}: no tension line found")

    # Slide 5 must have exactly 3 outcomes
    outcomes = copy.get('slide5', {}).get('outcomes', [])
    if len(outcomes) != 3:
        errors.append(f"slide5: must have exactly 3 outcomes, got {len(outcomes)}")
    for outcome in outcomes:
        low = outcome.lower()
        if any(word in low for word in ("step ", "first,", "second,", "third,", "process", "workflow")):
            errors.append("slide5: outcomes must be outcome lines, not process steps")

    # Every section needs a headline
    for key in ('slide1', 'slide2', 'slide3', 'slide4', 'slide5', 'slide6'):
        if not copy.get(key, {}).get('headline', '').strip():
            errors.append(f"{key}: missing headline")

    # slide6 headline should be a strong, short risk-reversal (not enforced as exact string)
    s6_headline = copy.get('slide6', {}).get('headline', '').strip()
    if len(s6_headline) > 60:
        errors.append("slide6: headline is too long — keep it under 60 characters")

    # Pricing must mention a monthly figure
    pricing = copy.get('slide6', {}).get('pricing', '')
    if not pricing.strip():
        errors.append("slide6: pricing line is missing")

    return errors


# ---------------------------------------------------------------------------
# Variable extraction
# ---------------------------------------------------------------------------

def _infer_acquisition(p: dict) -> str:
    ad   = (p.get('ad_status') or '').lower()
    outb = (p.get('outbound_status') or '').lower()
    hire = (p.get('hiring_signal') or '').lower()
    if outb == 'active_outbound':
        return "active outbound"
    if 'sdr' in hire:
        return "early SDR hire, no outbound infrastructure"
    if ad == 'running_ads' and outb == 'no_outbound':
        return "paid ads only, no cold outbound"
    if ad == 'running_ads':
        return "paid ads and inbound"
    return "inbound and referrals"


def _extract_pain(notes: str) -> str:
    if not notes:
        return ""
    m = re.search(r"Pain Point:\s*(.*)", notes)
    if not m:
        return ""
    raw = m.group(1).strip().rstrip('.')
    core_m = re.search(r'\b(?:on|with|around|from)\s+(.+)$', raw, re.IGNORECASE)
    return core_m.group(1).lower() if core_m else raw


def _extract_variables(p: dict) -> dict:
    comps = [c.strip() for c in (p.get('competitors') or '').split(',') if c.strip()]
    return {
        "company_name":        (p.get('company') or '').strip(),
        "contact_name":        ((p.get('name') or 'Founder').split()[0]),
        "industry":            (p.get('niche') or '').strip(),
        "target_buyer":        (p.get('icp') or '').strip(),
        "core_offer":          (p.get('product_feature') or '').strip(),
        "current_acquisition": _infer_acquisition(p),
        "competitor_1":        comps[0] if len(comps) > 0 else "",
        "competitor_2":        comps[1] if len(comps) > 1 else "",
        "pain_point":          _extract_pain(p.get('notes') or ''),
        "date":                date.today().strftime('%B %Y'),
    }


# ---------------------------------------------------------------------------
# Template copy generation (no API key needed)
# ---------------------------------------------------------------------------

def _gen_copy_template(v: dict) -> dict:
    company = v['company_name']
    contact = v['contact_name']
    industry = v['industry']
    target  = v['target_buyer']
    acq     = v['current_acquisition']
    c1      = v['competitor_1']
    c2      = v['competitor_2']
    pain    = v['pain_point']
    dt      = v['date']

    target_short = target.split(',')[0].strip() if target else 'your buyers'
    comp_str = ' and '.join(filter(None, [c1, c2])) or 'competitors in your space'
    comp_verb = 'are' if (c1 and c2) else 'is'

    # Slide 1 headline — acquisition-specific
    if 'no cold' in acq:
        h1 = f"Why {company}\u2019s pipeline stops when the budget does"
    elif 'sdr' in acq.lower():
        h1 = f"Why {company}\u2019s outbound hire isn\u2019t producing pipeline yet"
    else:
        h1 = f"Why {company} is invisible to {target_short}"

    # Slide 2 headline
    if 'no cold' in acq:
        h2 = f"Paid traffic alone cannot build {company}\u2019s pipeline"
    elif 'sdr' in acq.lower():
        h2 = "One SDR hire is not an outbound system"
    else:
        h2 = f"{company}\u2019s growth depends on being found"

    pain_hook = pain or f"the core problem {company} solves"

    cards = [
        {
            "title": f"{target_short} never see {company}",
            "body":  f"Without outbound, {company} only reaches buyers who find them first. "
                     f"The rest of the market is untouched.",
        },
        {
            "title": "Spend \u00a31 to acquire every lead",
            "body":  "Remove the budget and the pipeline stops. Not gradually \u2014 immediately. "
                     "You are renting your pipeline, not building it.",
        },
        {
            "title": f"{comp_str} {comp_verb} already in those inboxes",
            "body":  f"They reached {target_short} first. "
                     f"That conversation happened without {company}.",
        },
        {
            "title": "CPAs compound as audiences saturate",
            "body":  "What works at \u00a330 per lead becomes \u00a390 within 18 months. "
                     "The ceiling drops as the audience shrinks.",
        },
    ]

    steps = [
        f"Identify {target_short} before they start searching for a solution",
        f"Research every contact \u2014 {pain_hook[:55]} as the hook",
        "Reach them with sequences that read like a peer sent them",
        "Classify replies the same day \u2014 interested leads move to a booked call",
    ]

    terminal = [
        f'$ claude run leadgen-agent --company "{company}"',
        f'> Scanning {industry or "target market"} for {target_short}...',
        f'> Pain point identified: {pain_hook[:48]}',
        f'> 47 qualified contacts found in target segment',
        f'> Sequences queued: 47  |  Expected replies: 8\u201312',
        f'Done. {company} pipeline ready.',
    ]

    h4 = f"{comp_str} {comp_verb} already in your buyers\u2019 inboxes"

    outcomes = [
        f"We find the {target_short} already dealing with {pain_hook} "
        f"\u2014 before they go looking for a solution.",
        "We reach them directly. Every message is grounded in something real about their business.",
        f"We turn replies into booked conversations with {company}\u2019s team. Nothing is missed.",
    ]

    return {
        "slide1": {
            "headline":    h1,
            "subline":     f"{industry}  |  Outbound Pipeline Analysis" if industry
                           else "Outbound Pipeline Analysis",
            "stats":       [
                {"number": "4\u20138", "label": "qualified calls/month"},
                {"number": "30",       "label": "day pilot"},
                {"number": "\u00a30",  "label": "if no results"},
            ],
            "prepared_for": f"Prepared for {contact} at {company}  \u00b7  {dt}",
            "tension": "Pipeline should not stop when the budget does.",
        },
        "slide2": {
            "headline": h2,
            "cards":    cards,
            "tension":  "This caps your pipeline. Hard.",
        },
        "slide3": {
            "headline":       "Powered by Claude Code",
            "steps":          steps,
            "terminal_lines": terminal,
            "tension":        "First contact happens before inbound does.",
        },
        "slide4": {
            "headline": h4,
            "comp1": {
                "name":  c1 or "Competitor A",
                "line1": f"Running structured cold outreach into {target_short}",
                "line2": f"Booking first conversations before {company} enters the picture",
            },
            "comp2": {
                "name":  c2 or "Competitor B",
                "line1": f"Reaching {industry or 'your market'} buyers at scale",
                "line2": "Capturing demand before it becomes a search",
            },
            "insight": f"They get first contact. You do not.\n"
                       f"First contact sets the frame. "
                       f"That matters more than having a better product.",
            "tension": "They get first contact. You do not.",
        },
        "slide5": {
            "headline": "The System",
            "outcomes": outcomes,
            "tension": "Before they start searching is when this works.",
        },
        "slide6": {
            "headline": "No results. No charge.",
            "metrics":  [
                {"number": "4\u20138", "label": "qualified calls per month"},
                {"number": "30",       "label": "days to first results"},
                {"number": "\u00a30",  "label": "cost if no qualified conversations"},
            ],
            "cta_body": (
                f"Book a 20-minute call. We will walk through the target list, "
                f"the sequence, and expected volume \u2014 specific to {company}."
            ),
            "pricing": "Pilot: Free  |  Ongoing: from \u00a31,500/month  |  No long-term contract",
            "tension": "If the pipeline stops, you do not pay.",
        },
    }


# ---------------------------------------------------------------------------
# AI copy generation
# ---------------------------------------------------------------------------

def _gen_copy_ai(v: dict) -> dict:
    """Call Claude to generate bespoke slide copy. Returns copy dict."""
    import anthropic

    client = anthropic.Anthropic()
    user_msg = "Generate the deck copy now. Return only the JSON object."

    resp = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=3000,
        system=build_deck_system_prompt(v),
        messages=[{"role": "user", "content": user_msg}],
    )

    text = resp.content[0].text.strip()
    # Strip markdown fences if Claude wraps the JSON
    text = re.sub(r'^```[a-z]*\n?', '', text)
    text = re.sub(r'\n?```$', '', text.strip())
    return json.loads(text)


# ---------------------------------------------------------------------------
# Slide renderers
# ---------------------------------------------------------------------------

def _cover_badge(slide, x: float, y: float):
    """ANTIGRAVITY brand badge."""
    badge = slide.shapes.add_shape(1, _I(x), _I(y), _I(1.55), _I(0.26))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(ACCENTS[0])
    badge.line.fill.background()
    tf_b = badge.text_frame
    tf_b.margin_left = _I(0.09)
    tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
    p = tf_b.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "ANTIGRAVITY"
    r.font.size = _P(7)
    r.font.bold = True
    r.font.color.rgb = _rgb("FFFFFF")


def _slide1_cover(prs: Presentation, v: dict, c: dict):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[0])
    _footer(slide, 1)

    layout     = STYLE.get("cover_layout", "hero_left")
    h1_size    = STYLE.get("h1", 44)
    company    = (v.get('company_name') or 'A')[0].upper()
    stats      = c['slide1']['stats']
    x_off      = 0.55 if STYLE.get("accent_bar") == "left" else 0.45

    # ── Layout 1: hero_left (dark_indigo) ──────────────────────────────────
    # Headline fills left 2/3 · stat cards stacked on right
    if layout == "hero_left":
        ghost = slide.shapes.add_textbox(_I(7.5), _I(0.5), _I(6.0), _I(6.0))
        tf_g  = ghost.text_frame
        tf_g.margin_left = tf_g.margin_right = tf_g.margin_top = tf_g.margin_bottom = 0
        p_g = tf_g.paragraphs[0]
        p_g.alignment = PP_ALIGN.RIGHT
        r_g = p_g.add_run()
        r_g.text = company
        r_g.font.size = _P(400)
        r_g.font.bold = True
        r_g.font.color.rgb = _rgb(CARD_BG)

        _cover_badge(slide, x_off, 0.17)
        tf_h = _tx(slide, x_off, 0.72, 8.6, 2.8)
        _run(tf_h, c['slide1']['headline'], h1_size, TEXT_PRI, bold=True)
        tf_sub = _tx(slide, x_off, 3.55, 8.0, 0.45)
        _run(tf_sub, c['slide1']['subline'], 12, TEXT_SEC)
        tf_pre = _tx(slide, x_off, 6.82, 9.0, 0.28)
        _run(tf_pre, c['slide1']['prepared_for'], 9, TEXT_MUT)

        card_w, card_h = 3.7, 1.5
        for i, (stat, y) in enumerate(zip(stats, [1.4, 3.05, 4.7])):
            _card(slide, 9.18, y, card_w, card_h)
            tf_n = _tx(slide, 9.4, y + 0.1, card_w - 0.44, 0.88)
            _run(tf_n, stat['number'], 46, ACCENTS[i % len(ACCENTS)], bold=True)
            tf_l = _tx(slide, 9.4, y + 0.95, card_w - 0.44, 0.42)
            _run(tf_l, stat['label'], 10, TEXT_SEC)

    # ── Layout 2: centered_bottom (charcoal_gold) ───────────────────────────
    # Headline centered · stat cards horizontal row at bottom
    elif layout == "centered_bottom":
        # Faint company word behind headline
        ghost = slide.shapes.add_textbox(_I(0.3), _I(0.8), _I(13.0), _I(4.5))
        tf_g  = ghost.text_frame
        tf_g.margin_left = tf_g.margin_right = tf_g.margin_top = tf_g.margin_bottom = 0
        p_g   = tf_g.paragraphs[0]
        p_g.alignment = PP_ALIGN.CENTER
        r_g   = p_g.add_run()
        r_g.text = (v.get('company_name') or '').upper()
        r_g.font.size   = _P(120)
        r_g.font.bold   = True
        r_g.font.color.rgb = _rgb(CARD_BG)

        _cover_badge(slide, 5.89, 0.17)   # centered badge
        tf_h = _tx(slide, 0.6, 1.45, 12.13, 2.4)
        _run(tf_h, c['slide1']['headline'], h1_size, TEXT_PRI, bold=True, align=PP_ALIGN.CENTER)
        tf_sub = _tx(slide, 0.6, 3.95, 12.13, 0.45)
        _run(tf_sub, c['slide1']['subline'], 12, TEXT_SEC, align=PP_ALIGN.CENTER)
        tf_pre = _tx(slide, x_off, 6.82, 12.43, 0.28)
        _run(tf_pre, c['slide1']['prepared_for'], 9, TEXT_MUT, align=PP_ALIGN.CENTER)

        # Horizontal stat row
        card_w, card_h = 3.88, 1.55
        for i, (stat, sx) in enumerate(zip(stats, [0.45, 4.57, 8.69])):
            _card(slide, sx, 4.65, card_w, card_h)
            _card_stripe(slide, sx, 4.65, card_w, card_h, ACCENTS[i % len(ACCENTS)])
            tf_n = _tx(slide, sx + 0.22, 4.75, card_w - 0.44, 0.85)
            _run(tf_n, stat['number'], 42, ACCENTS[i % len(ACCENTS)], bold=True)
            tf_l = _tx(slide, sx + 0.22, 5.55, card_w - 0.44, 0.42)
            _run(tf_l, stat['label'], 10, TEXT_SEC)

    # ── Layout 3: reversed_split (midnight_teal) ────────────────────────────
    # Stats stacked LEFT · headline on RIGHT · bottom accent bar
    elif layout == "reversed_split":
        card_w, card_h = 3.9, 1.55
        for i, (stat, sy) in enumerate(zip(stats, [0.95, 2.7, 4.45])):
            _card(slide, 0.35, sy, card_w, card_h)
            tf_n = _tx(slide, 0.57, sy + 0.1, card_w - 0.44, 0.85)
            _run(tf_n, stat['number'], 40, ACCENTS[i % len(ACCENTS)], bold=True)
            tf_l = _tx(slide, 0.57, sy + 0.95, card_w - 0.44, 0.42)
            _run(tf_l, stat['label'], 10, TEXT_SEC)

        # Vertical divider line
        div = slide.shapes.add_shape(1, _I(4.6), _I(0.55), _I(0.02), _I(6.3))
        div.fill.solid()
        div.fill.fore_color.rgb = _rgb(CARD_BD)
        div.line.fill.background()

        _cover_badge(slide, 4.95, 0.22)
        tf_h = _tx(slide, 4.95, 0.75, 8.0, 3.2)
        _run(tf_h, c['slide1']['headline'], h1_size, TEXT_PRI, bold=True)
        tf_sub = _tx(slide, 4.95, 4.0, 8.0, 0.55)
        _run(tf_sub, c['slide1']['subline'], 12, TEXT_SEC)
        tf_pre = _tx(slide, 4.95, 6.82, 8.0, 0.28)
        _run(tf_pre, c['slide1']['prepared_for'], 9, TEXT_MUT)

    # ── Layout 4: fullbleed_bold (dark_slate_red) ───────────────────────────
    # Very thick accent bar · huge ghost company initial · oversized headline
    # · stats in a horizontal row below headline
    elif layout == "fullbleed_bold":
        # Very visible ghost initial
        ghost = slide.shapes.add_textbox(_I(6.5), _I(0.2), _I(6.8), _I(7.0))
        tf_g  = ghost.text_frame
        tf_g.margin_left = tf_g.margin_right = tf_g.margin_top = tf_g.margin_bottom = 0
        p_g   = tf_g.paragraphs[0]
        p_g.alignment = PP_ALIGN.RIGHT
        r_g   = p_g.add_run()
        r_g.text = company
        r_g.font.size  = _P(520)
        r_g.font.bold  = True
        r_g.font.color.rgb = _rgb("1F1010")

        # Company supertitle in accent
        tf_sup = _tx(slide, x_off, 0.28, 9.0, 0.35)
        _run(tf_sup, (v.get('company_name') or '').upper(), 10,
             ACCENTS[0], bold=True)

        _cover_badge(slide, x_off, 0.68)
        tf_h = _tx(slide, x_off, 1.1, 11.5, 3.0)
        _run(tf_h, c['slide1']['headline'], h1_size, TEXT_PRI, bold=True)
        tf_sub = _tx(slide, x_off, 4.15, 10.0, 0.45)
        _run(tf_sub, c['slide1']['subline'], 12, TEXT_SEC)
        tf_pre = _tx(slide, x_off, 6.82, 9.0, 0.28)
        _run(tf_pre, c['slide1']['prepared_for'], 9, TEXT_MUT)

        card_w, card_h = 3.88, 1.35
        for i, (stat, sx) in enumerate(zip(stats, [0.45, 4.57, 8.69])):
            _card(slide, sx, 4.72, card_w, card_h)
            _card_stripe(slide, sx, 4.72, card_w, card_h, ACCENTS[i % len(ACCENTS)])
            tf_n = _tx(slide, sx + 0.22, 4.82, card_w - 0.44, 0.75)
            _run(tf_n, stat['number'], 38, ACCENTS[i % len(ACCENTS)], bold=True)
            tf_l = _tx(slide, sx + 0.22, 5.5, card_w - 0.44, 0.42)
            _run(tf_l, stat['label'], 9, TEXT_SEC)


def _slide2_problem(prs: Presentation, v: dict, c: dict):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[1])
    _section_label(slide, "The Problem", ACCENTS[1])
    _footer(slide, 2)

    h2    = STYLE.get("h2", 34)
    body  = STYLE.get("body", 11)
    x_off = 0.55 if STYLE.get("accent_bar") == "left" else 0.45
    stripe_mode = STYLE.get("card_stripe", "dot")

    tf_h = _tx(slide, x_off, 0.38, 12.4, 0.9)
    _run(tf_h, c['slide2']['headline'], h2, TEXT_PRI, bold=True)

    card_w = 6.1
    card_h = 2.2
    xs = [x_off, x_off + 6.3]
    ys = [1.45, 3.85]
    cards = c['slide2']['cards'][:4]
    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            if idx >= len(cards):
                break
            cx, cy = xs[col], ys[row]
            _card(slide, cx, cy, card_w, card_h)
            accent_c = ACCENTS[(idx + 1) % len(ACCENTS)]
            _card_stripe(slide, cx, cy, card_w, card_h, accent_c)
            if stripe_mode == "dot":
                dot = slide.shapes.add_shape(1, _I(cx + 0.22), _I(cy + 0.22),
                                             _I(0.08), _I(0.08))
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(accent_c)
                dot.line.fill.background()
                tf_t = _tx(slide, cx + 0.42, cy + 0.16, card_w - 0.6, 0.4)
            elif stripe_mode == "top":
                tf_t = _tx(slide, cx + 0.22, cy + 0.22, card_w - 0.44, 0.4)
            elif stripe_mode == "left":
                tf_t = _tx(slide, cx + 0.28, cy + 0.16, card_w - 0.5, 0.4)
            else:
                tf_t = _tx(slide, cx + 0.22, cy + 0.16, card_w - 0.44, 0.4)
            _run(tf_t, cards[idx]['title'], STYLE.get("h3", 16), TEXT_PRI, bold=True)
            tf_b2 = _tx(slide, cx + 0.22, cy + 0.6, card_w - 0.44, 1.45)
            _run(tf_b2, cards[idx]['body'], body, TEXT_SEC)

    tf_t = _tx(slide, x_off, 6.2, 12.4, 0.5)
    _run(tf_t, c['slide2']['tension'], 16, DANGER, bold=True)


def _slide3_claude(prs: Presentation, v: dict, c: dict):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[2])
    _section_label(slide, "How It Works", ACCENTS[2])
    _footer(slide, 3)

    h2    = STYLE.get("h2", 34)
    body  = STYLE.get("body", 11)
    x_off = 0.55 if STYLE.get("accent_bar") == "left" else 0.45
    mono  = STYLE.get("section_mono", False)

    tf_h = _tx(slide, x_off, 0.38, 12.4, 0.75)
    _run(tf_h, c['slide3'].get('headline', "The agent that builds your pipeline"), h2, TEXT_PRI, bold=True)

    steps  = c['slide3']['steps']
    step_y = 1.38
    for i, step in enumerate(steps[:4]):
        if mono:
            # Outlined badge for teal theme
            badge = slide.shapes.add_shape(1, _I(x_off), _I(step_y), _I(0.46), _I(0.46))
            badge.fill.background()
            badge.line.color.rgb = _rgb(ACCENTS[2])
            badge.line.width = _P(1.5)
        else:
            badge = slide.shapes.add_shape(1, _I(x_off), _I(step_y), _I(0.46), _I(0.46))
            badge.fill.solid()
            badge.fill.fore_color.rgb = _rgb(ACCENTS[2])
            badge.line.fill.background()
        tf_n = badge.text_frame
        tf_n.margin_left = tf_n.margin_right = tf_n.margin_top = tf_n.margin_bottom = 0
        pn = tf_n.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run()
        rn.text = str(i + 1)
        rn.font.size = _P(12)
        rn.font.bold = True
        rn.font.color.rgb = _rgb(ACCENTS[2] if mono else "FFFFFF")
        if mono:
            rn.font.name = "Courier New"

        tf_s = _tx(slide, x_off + 0.63, step_y + 0.04, 5.25, 0.55)
        _run(tf_s, step, body + 1, TEXT_PRI, mono=mono)
        step_y += 1.35

    # Right — terminal card
    _card(slide, 6.8, 1.3, 6.1, 5.55)

    # Terminal header bar
    bar = slide.shapes.add_shape(
        1, _I(6.8), _I(1.3), _I(6.1), _I(0.36)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb("1C2333")
    bar.line.fill.background()

    # Terminal title text in header
    tf_title = slide.shapes.add_textbox(_I(7.55), _I(1.33), _I(4.5), _I(0.3))
    tf_title = tf_title.text_frame
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    pt = tf_title.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    rt = pt.add_run()
    rt.text = "agent output"
    rt.font.size = _P(8)
    rt.font.color.rgb = _rgb(TEXT_MUT)
    rt.font.name = "Courier New"

    # Traffic lights
    for xi, col in enumerate(["EF4444", "F59E0B", "10B981"]):
        dot = slide.shapes.add_shape(
            1,
            _I(7.03 + xi * 0.24), _I(1.42),
            _I(0.13), _I(0.13),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(col)
        dot.line.fill.background()

    # Terminal lines — slightly larger font
    term_lines = c['slide3']['terminal_lines']
    tf_term = _tx(slide, 7.05, 1.82, 5.65, 4.8)
    for i, line in enumerate(term_lines):
        is_cmd = line.startswith('$')
        color  = TERM_CMD if is_cmd else (TERM_GRN if line.startswith('>') else TEXT_SEC)
        _run(tf_term, line, 9.5, color, mono=True, space_before=(5 if i else 0))


def _slide4_market(prs: Presentation, v: dict, c: dict):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[3])
    _section_label(slide, "Market Reality", ACCENTS[3])
    _footer(slide, 4)

    h2    = STYLE.get("h2", 34)
    body  = STYLE.get("body", 11)
    x_off = 0.55 if STYLE.get("accent_bar") == "left" else 0.45

    tf_h = _tx(slide, x_off, 0.38, 12.4, 0.85)
    _run(tf_h, c['slide4']['headline'], h2 - 2, TEXT_PRI, bold=True)

    comp_card_w = 5.9
    comp_card_h = 3.55
    for i, (comp_key, cx) in enumerate([('comp1', x_off), ('comp2', x_off + 6.1)]):
        comp   = c['slide4'][comp_key]
        accent = ACCENTS[i]
        _card(slide, cx, 1.4, comp_card_w, comp_card_h)
        _card_stripe(slide, cx, 1.4, comp_card_w, comp_card_h, accent)

        text_x = cx + (0.22 if STYLE.get("card_stripe") == "left" else 0.25)
        tf_name = _tx(slide, text_x, 1.63, comp_card_w - 0.5, 0.45)
        _run(tf_name, comp['name'], STYLE.get("h3", 16), TEXT_PRI, bold=True)
        tf_l1 = _tx(slide, text_x, 2.18, comp_card_w - 0.5, 0.52)
        _run(tf_l1, comp['line1'], body, TEXT_SEC)
        tf_l2 = _tx(slide, text_x, 2.82, comp_card_w - 0.5, 1.0)
        _run(tf_l2, comp['line2'], body, TEXT_MUT)

    _card(slide, x_off, 5.18, 12.43, 1.55)
    stripe2 = slide.shapes.add_shape(1, _I(x_off), _I(5.18), _I(0.07), _I(1.55))
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = _rgb(DANGER)
    stripe2.line.fill.background()

    lines = c['slide4']['insight'].split('\n')
    tf_ins = _tx(slide, x_off + 0.3, 5.34, 12.0, 1.25)
    _run(tf_ins, lines[0], 16, DANGER, bold=True)
    if len(lines) > 1:
        _run(tf_ins, lines[1], body, TEXT_SEC, space_before=6)


def _slide5_system(prs: Presentation, v: dict, c: dict):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[0])
    _section_label(slide, "The System", ACCENTS[0])
    _footer(slide, 5)

    h2       = STYLE.get("h2", 34)
    body     = STYLE.get("body", 11)
    x_off    = 0.55 if STYLE.get("accent_bar") == "left" else 0.45
    s5_layout = STYLE.get("slide5_layout", "stripe_cards")
    outcomes = c['slide5']['outcomes']

    tf_h = _tx(slide, x_off, 0.36, 12.4, 0.75)
    _run(tf_h, c['slide5']['headline'], h2, TEXT_PRI, bold=True)

    if s5_layout == "pipe_lines":
        # Minimal: 3 full-width lines separated by thin rules, no cards
        accent_cycle = [ACCENTS[1], ACCENTS[2], ACCENTS[3]]
        ys = [1.42, 3.1, 4.78]
        for i, (outcome, oy) in enumerate(zip(outcomes[:3], ys)):
            accent = accent_cycle[i]
            # Pipe glyph in accent
            tf_pipe = _tx(slide, x_off, oy, 0.5, 1.4)
            _run(tf_pipe, "|", 48, accent, bold=True, mono=True)
            tf_o = _tx(slide, x_off + 0.55, oy + 0.2, 11.8, 1.1)
            _run(tf_o, outcome, body + 3, TEXT_PRI)
            # Thin separator rule
            if i < 2:
                rule = slide.shapes.add_shape(1, _I(x_off), _I(oy + 1.5),
                                              _I(12.43), _I(0.01))
                rule.fill.solid()
                rule.fill.fore_color.rgb = _rgb(CARD_BD)
                rule.line.fill.background()
    else:
        # Default: stripe cards
        card_w = 12.43
        card_h = 1.6
        ys = [1.32, 3.1, 4.88]
        accent_cycle = [ACCENTS[1], ACCENTS[2], ACCENTS[3]]
        for i, (outcome, cy) in enumerate(zip(outcomes[:3], ys)):
            accent = accent_cycle[i]
            _card(slide, x_off, cy, card_w, card_h)
            stripe = slide.shapes.add_shape(
                1, _I(x_off), _I(cy), _I(0.07), _I(card_h))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = _rgb(accent)
            stripe.line.fill.background()
            tf_n = _tx(slide, x_off + 0.29, cy + 0.38, 0.68, 0.7)
            _run(tf_n, f"0{i+1}", 28, accent, bold=True)
            tf_o = _tx(slide, x_off + 1.07, cy + 0.22, card_w - 1.25, 1.2)
            _run(tf_o, outcome, body + 2, TEXT_PRI)


def _slide6_cta(prs: Presentation, v: dict, c: dict):  # noqa: ARG001
    slide = _blank_slide(prs)
    _set_bg(slide)
    _accent_bar(slide, ACCENTS[3])
    _section_label(slide, "Pilot + CTA", ACCENTS[3])
    _footer(slide, 6)

    h1    = STYLE.get("h1", 44)
    body  = STYLE.get("body", 11)
    x_off = 0.55 if STYLE.get("accent_bar") == "left" else 0.45

    tf_h = _tx(slide, x_off, 0.36, 12.4, 0.82)
    _run(tf_h, c['slide6']['headline'], h1, DANGER, bold=True)

    metrics   = c['slide6']['metrics']
    metric_w  = 3.88
    metric_h  = 1.55
    metric_xs = [x_off, x_off + 4.12, x_off + 8.24]
    for i, (metric, mx) in enumerate(zip(metrics[:3], metric_xs)):
        _card(slide, mx, 1.32, metric_w, metric_h)
        _card_stripe(slide, mx, 1.32, metric_w, metric_h, ACCENTS[i % len(ACCENTS)])
        tf_num = _tx(slide, mx + 0.22, 1.42, metric_w - 0.44, 0.88)
        _run(tf_num, metric['number'], 46, ACCENTS[i % len(ACCENTS)], bold=True)
        tf_lbl = _tx(slide, mx + 0.22, 2.18, metric_w - 0.44, 0.55)
        _run(tf_lbl, metric['label'], 10, TEXT_SEC)

    _card(slide, x_off, 3.1, 12.43, 2.8)

    tf_ctab = _tx(slide, x_off + 0.33, 3.26, 11.85, 0.62)
    _run(tf_ctab, c['slide6']['cta_body'], body + 1, TEXT_SEC)

    btn = slide.shapes.add_shape(1, _I(x_off + 0.33), _I(4.08), _I(5.2), _I(0.58))
    btn.fill.solid()
    btn.fill.fore_color.rgb = _rgb(ACCENTS[0])
    btn.line.fill.background()
    _shadow(btn)
    tf_btn = btn.text_frame
    tf_btn.margin_left = tf_btn.margin_right = tf_btn.margin_top = tf_btn.margin_bottom = 0
    pb = tf_btn.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = f"Book a Call  \u2192  {CTA_LINK}"
    rb.font.size = _P(12)
    rb.font.bold = True
    rb.font.color.rgb = _rgb("FFFFFF")

    tf_pr = _tx(slide, x_off + 0.33, 5.55, 12.0, 0.3)
    _run(tf_pr, c['slide6']['pricing'], 9.5, TEXT_MUT, align=PP_ALIGN.CENTER)


def _convert_deck_to_pdf(filepath: str) -> tuple[str | None, str]:
    """Best-effort PPTX to PDF conversion using LibreOffice when available."""
    soffice = shutil.which("soffice")
    if not soffice:
        return _convert_deck_to_pdf_windows(filepath)

    outdir = os.path.dirname(os.path.abspath(filepath))
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", filepath, "--outdir", outdir],
            capture_output=True,
            text=True,
            timeout=120,
            check=True,
        )
    except Exception as exc:
        return None, f"failed: {exc}"

    pdf_path = os.path.splitext(filepath)[0] + ".pdf"
    if not os.path.exists(pdf_path):
        detail = result.stderr.strip() or result.stdout.strip() or "no pdf produced"
        return None, f"failed: {detail}"
    return pdf_path, "ok"


def _convert_deck_to_pdf_windows(filepath: str) -> tuple[str | None, str]:
    """Windows fallback: export PDF through PowerPoint COM via PowerShell."""
    if os.name != "nt":
        return None, "skipped: soffice not installed"

    powershell = shutil.which("powershell")
    if not powershell:
        return None, "skipped: powershell not available"

    pptx_path = os.path.abspath(filepath)
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    ps_pptx = pptx_path.replace("'", "''")
    ps_pdf = pdf_path.replace("'", "''")
    command = (
        "$ErrorActionPreference = 'Stop'; "
        "$ppt = $null; $deck = $null; "
        "try { "
        "$ppt = New-Object -ComObject PowerPoint.Application; "
        "$ppt.Visible = -1; "
        f"$deck = $ppt.Presentations.Open('{ps_pptx}', $false, $false, $false); "
        f"$deck.SaveAs('{ps_pdf}', 32); "
        "} finally { "
        "if ($deck) { $deck.Close() } "
        "if ($ppt) { $ppt.Quit() } "
        "}"
    )
    try:
        result = subprocess.run(
            [powershell, "-NoProfile", "-Command", command],
            capture_output=True,
            text=True,
            timeout=120,
            check=True,
        )
    except Exception as exc:
        stderr = ""
        stdout = ""
        if hasattr(exc, "stderr") and exc.stderr:
            stderr = str(exc.stderr).strip()
        if hasattr(exc, "stdout") and exc.stdout:
            stdout = str(exc.stdout).strip()
        detail = stderr or stdout or str(exc)
        return None, f"failed: {detail}"

    if not os.path.exists(pdf_path):
        detail = result.stderr.strip() or result.stdout.strip() or "PowerPoint export did not create a PDF"
        return None, f"failed: {detail}"
    return pdf_path, "ok (PowerPoint)"


def _rasterize_pdf(pdf_path: str) -> tuple[list[str], str]:
    """Best-effort PDF rasterization using pdftoppm when available."""
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return [], "skipped: pdftoppm not installed"

    out_prefix = os.path.splitext(pdf_path)[0] + "_slide"
    try:
        subprocess.run(
            [pdftoppm, "-jpeg", "-r", "150", pdf_path, out_prefix],
            capture_output=True,
            text=True,
            timeout=120,
            check=True,
        )
    except Exception as exc:
        return [], f"failed: {exc}"

    directory = os.path.dirname(out_prefix) or "."
    basename = os.path.basename(out_prefix)
    images = sorted(
        os.path.join(directory, name)
        for name in os.listdir(directory)
        if name.startswith(basename) and name.endswith(".jpg")
    )
    return images, "ok" if images else "failed: no images produced"


def run_deck_qa(copy: dict, variables: dict, filepath: str | None = None) -> dict:
    """Run copy checks and optional file-based QA steps for a generated deck."""
    issues = _validate_copy(copy)
    company = variables.get("company_name", "")

    for key in ("slide1", "slide2", "slide3", "slide4", "slide5", "slide6"):
        slide_text = json.dumps(copy.get(key, {}))
        if company and key != "slide4" and company not in slide_text:
            issues.append(f"{key}: missing explicit company reference")

    terminal = " ".join(copy.get("slide3", {}).get("terminal_lines", []))
    target_buyer = variables.get("target_buyer", "")
    target_check = target_buyer.split(",")[0].strip() if target_buyer else ""
    pain_point = variables.get("pain_point", "")
    pain_check = pain_point.split(",")[0].strip() if pain_point else ""
    required_terminal_values = (
        variables.get("company_name", ""),
        target_check,
        pain_check,
    )
    for required in required_terminal_values:
        if required and required not in terminal:
            issues.append("slide3: terminal output must reference company, target buyer, and pain point")
            break

    slide4_text = json.dumps(copy.get("slide4", {}))
    for competitor in (variables.get("competitor_1", ""), variables.get("competitor_2", "")):
        if competitor and competitor not in slide4_text:
            issues.append(f"slide4: competitor '{competitor}' missing from market reality copy")

    qa_report = {
        "copy_issues": issues,
        "pdf_status": "not_run",
        "raster_status": "not_run",
        "pdf_path": None,
        "slide_images": [],
    }

    if filepath:
        pdf_path, pdf_status = _convert_deck_to_pdf(filepath)
        qa_report["pdf_status"] = pdf_status
        qa_report["pdf_path"] = pdf_path
        if pdf_path:
            images, raster_status = _rasterize_pdf(pdf_path)
            qa_report["raster_status"] = raster_status
            qa_report["slide_images"] = images

    return qa_report


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_deck(prospect: dict, prefer_pdf: bool = True, theme: str = "dark_indigo") -> str:
    """
    Generate a bespoke pitch deck for a prospect.

    Raises ValueError on validation failure.
    Returns filepath to the generated PDF when available, otherwise PPTX if
    prefer_pdf is False.

    theme: one of "dark_indigo", "charcoal_gold", "midnight_teal", "dark_slate_red"
    """
    company = (prospect.get('company') or '').strip()
    if not company:
        raise ValueError("prospect must have a 'company' field")

    has_enrichment = any(
        (prospect.get(f) or '').strip()
        for f in ('niche', 'icp', 'website_headline', 'notes')
    )
    if not has_enrichment:
        raise ValueError(
            f"'{company}' has no enrichment. Run research_agent first."
        )

    v = _extract_variables(prospect)

    # Choose copy source
    has_api_key = bool(os.getenv("ANTHROPIC_API_KEY", "").strip())
    if has_api_key:
        try:
            copy = _gen_copy_ai(v)
            # Validate immediately; fall back to template if AI copy fails QA
            ai_issues = _validate_copy(copy)
            if ai_issues:
                copy = _gen_copy_template(v)
                copy["_mode"] = f"template (AI copy failed QA: {'; '.join(ai_issues[:2])})"
        except Exception as exc:
            copy = _gen_copy_template(v)
            copy["_mode"] = f"template (AI failed: {exc})"
    else:
        copy = _gen_copy_template(v)
        copy["_mode"] = "template"

    # Render — apply theme while building slides so all globals are swapped
    with _apply_theme(theme):
        prs = Presentation()
        prs.slide_width  = _I(13.33)   # scaled: 13.33 in (full) or 10 in (compact)
        prs.slide_height = _I(7.5)     # scaled: 7.5 in  (full) or 5.625 in (compact)

        _slide1_cover(prs, v, copy)
        _slide2_problem(prs, v, copy)
        _slide3_claude(prs, v, copy)
        _slide4_market(prs, v, copy)
        _slide5_system(prs, v, copy)
        _slide6_cta(prs, v, copy)

    safe = "".join(c if c.isalnum() else "_" for c in company).lower()
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    filepath = os.path.join(OUTPUT_DIR, f"deck_{safe}.pptx")
    prs.save(filepath)

    qa_report = run_deck_qa(copy, v, filepath)
    if qa_report["copy_issues"]:
        print(f"  [WARN] {filepath} copy QA: {'; '.join(qa_report['copy_issues'][:3])}")

    mode = copy.get("_mode", "ai" if has_api_key else "template")
    print(f"  [OK] {filepath}  ({mode}; pdf={qa_report['pdf_status']}; raster={qa_report['raster_status']})")

    if prefer_pdf:
        pdf_path = qa_report.get("pdf_path")
        if pdf_path and os.path.exists(pdf_path):
            return pdf_path
        # PDF export failed (LibreOffice not installed, etc.) — return PPTX
        print(f"  [INFO] PDF export failed ({qa_report['pdf_status']}); returning PPTX")

    return filepath


# ---------------------------------------------------------------------------
# __main__ — test fixture (Apex Digital)
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    thin = {
        "name":    "James Cole",
        "company": "Apex Digital",
        "email":   "james@apexdigital.io",
        "status":  "qualified",
        "notes":   "",
    }

    rich = {
        "name":             "James Cole",
        "company":          "Apex Digital",
        "email":            "james@apexdigital.io",
        "status":           "qualified",
        "niche":            "B2B SaaS for construction project managers",
        "icp":              "mid-sized UK construction firms running 5+ concurrent projects",
        "website_headline": "Stop losing projects to miscommunication",
        "product_feature":  "real-time site-to-office sync with automated progress reporting",
        "competitors":      "Buildertrend, CoConstruct",
        "ad_status":        "running_ads",
        "outbound_status":  "no_outbound",
        "notable_result":   "Cut project overruns by 40% for Morrison Construction",
        "notes": (
            "[Research Hook]\n"
            "Pain Point: Construction PMs spend 6+ hours per week on manual status updates.\n"
            "Growth Signal: Recently rebranded with enterprise case studies added.\n"
            "Opener: Apex Digital positioning around miscommunication is sharp."
        ),
    }

    print("=" * 60)
    print("BEFORE -- thin data (should be rejected)")
    print("=" * 60)
    try:
        generate_deck(thin)
    except ValueError as e:
        print(f"  [OK] Correctly rejected: {e}")

    print()
    print("=" * 60)
    print("AFTER -- full enrichment")
    print("=" * 60)

    # Show copy preview
    v = _extract_variables(rich)
    copy = _gen_copy_template(v)

    print(f"\n  Slide 1 headline : {copy['slide1']['headline']}")
    print(f"  Slide 2 tension  : {copy['slide2']['tension']}")
    print(f"  Slide 4 insight  : {copy['slide4']['insight'].splitlines()[0]}")
    print(f"  Slide 5 outcome  : {copy['slide5']['outcomes'][0][:70]}...")

    print("\n  Tension lines per section:")
    print(f"    02 The Problem  : {copy['slide2']['tension']}")
    print(f"    04 Market       : {copy['slide4']['insight'].splitlines()[0]}")
    print(f"    05 System       : (outcome format — no hard pipeline stops)")

    print()
    try:
        path = generate_deck(rich)
        print(f"  Generated: {path}")
    except ValueError as e:
        print(f"  ERROR: {e}")
        sys.exit(1)
